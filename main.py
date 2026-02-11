# helo
import serial
import serial.tools.list_ports
import time
import os
import shutil
import wave
import difflib
import json
import sys
import threading
import socket
import speech_recognition as sr
from vosk import Model, KaldiRecognizer
import customtkinter as ctk 
import tkinter as tk
from tkinter import filedialog, scrolledtext, messagebox
import ctypes
from ctypes import wintypes

try:
    import docx
except ImportError:
    docx = None
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

BAUD = 921600                
TEMP_WAV = "temp_audio.wav"
MAX_GOOGLE_LATENCY = 2.0 

def text_to_int(text):
    text = text.lower().replace('-', ' ')
    
    ones = {
        'zero': 0, 'one': 1, 'won': 1, 'two': 2, 'to': 2, 'too': 2,
        'three': 3, 'four': 4, 'for': 4, 'five': 5, 'six': 6,
        'seven': 7, 'eight': 8, 'ate': 8, 'nine': 9, 'ten': 10
    }
    
    teens = {
        'eleven': 11, 'twelve': 12, 'thirteen': 13, 'fourteen': 14,
        'fifteen': 15, 'sixteen': 16, 'seventeen': 17, 'eighteen': 18, 'nineteen': 19
    }
    
    tens = {
        'twenty': 20, 'thirty': 30, 'forty': 40, 'fifty': 50,
        'sixty': 60, 'seventy': 70, 'eighty': 80, 'ninety': 90, 'hundred': 100
    }
    
    for word in text.split():
        if word.isdigit():
            num = int(word)
            if 1 <= num <= 100:
                return num
    
    words = text.split()
    total = 0
    current = 0
    
    for word in words:
        if word in ones:
            current += ones[word]
        elif word in teens:
            current += teens[word]
        elif word in tens:
            if tens[word] == 100:
                current = (current or 1) * 100
            else:
                current += tens[word]
        elif word == 'and':
            continue
    
    total += current
    
    if total == 0:
        for word in words:
            if word in ones:
                return ones[word]
            elif word in teens:
                return teens[word]
            elif word in tens:
                return teens[word]
    
    return total if 1 <= total <= 100 else None

def resource_path(relative_path):
    try:
        base_path = sys._MEIPASS
    except Exception:
        base_path = os.path.abspath(".")
    return os.path.join(base_path, relative_path)

class DROPFILES(ctypes.Structure):
    _fields_ = [
        ("pFiles", wintypes.DWORD),
        ("pt", wintypes.POINT),
        ("fNC", wintypes.BOOL),
        ("fWide", wintypes.BOOL),
    ]

def copy_files_to_clipboard(file_paths):
    user32 = ctypes.windll.user32
    kernel32 = ctypes.windll.kernel32

    user32.OpenClipboard.argtypes = [wintypes.HWND]
    user32.OpenClipboard.restype = wintypes.BOOL
    
    user32.EmptyClipboard.argtypes = []
    user32.EmptyClipboard.restype = wintypes.BOOL

    user32.CloseClipboard.argtypes = []
    user32.CloseClipboard.restype = wintypes.BOOL

    user32.SetClipboardData.argtypes = [wintypes.UINT, wintypes.HANDLE]
    user32.SetClipboardData.restype = wintypes.HANDLE

    kernel32.GlobalAlloc.argtypes = [wintypes.UINT, ctypes.c_size_t]
    kernel32.GlobalAlloc.restype = wintypes.HGLOBAL

    kernel32.GlobalLock.argtypes = [wintypes.HGLOBAL]
    kernel32.GlobalLock.restype = ctypes.c_void_p

    kernel32.GlobalUnlock.argtypes = [wintypes.HGLOBAL]
    kernel32.GlobalUnlock.restype = wintypes.BOOL
    
    kernel32.GlobalFree.argtypes = [wintypes.HGLOBAL]
    kernel32.GlobalFree.restype = wintypes.HGLOBAL

    offset = ctypes.sizeof(DROPFILES)
    length = sum(len(p) + 1 for p in file_paths) + 1
    size = offset + length * ctypes.sizeof(ctypes.c_wchar)

    buf = (ctypes.c_char * size)()
    df = DROPFILES.from_buffer(buf)
    df.pFiles = offset
    df.fWide = True
    df.fNC = False 
    df.pt = wintypes.POINT(0, 0) 

    for path in file_paths:
        array_t = ctypes.c_wchar * (len(path) + 1)
        path_buf = array_t.from_buffer(buf, offset)
        path_buf.value = path
        offset += ctypes.sizeof(array_t)

    buf[offset] = b'\00'
    buf[offset+1] = b'\00' 
    
    if not user32.OpenClipboard(None):
        return
        
    try:
        user32.EmptyClipboard()
        
        h_global = kernel32.GlobalAlloc(0x0042, size)
        if not h_global:
            return
            
        p_global = kernel32.GlobalLock(h_global)
        if not p_global:
            kernel32.GlobalFree(h_global)
            return
            
        ctypes.memmove(p_global, buf, size)
        kernel32.GlobalUnlock(h_global)
        
        user32.SetClipboardData(15, h_global) 
    except Exception as e:
        print(f"clipboard error: {e}")
    finally:
        user32.CloseClipboard()

def get_file_content(filepath):
    content = ""
    ext = os.path.splitext(filepath)[1].lower()
    try:
        if ext == '.txt':
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
        
        elif ext == '.docx' and docx:
            doc = docx.Document(filepath)
            content = " ".join([p.text for p in doc.paragraphs])
            
        elif ext == '.pdf' and PyPDF2:
            with open(filepath, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    content += page.extract_text() + " "
        
        elif (ext == '.pptx' or ext == '.ppt') and Presentation:
            prs = Presentation(filepath)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + " "
    except: pass
    return content.lower()

ctk.set_appearance_mode("Dark")
ctk.set_default_color_theme("dark-blue")

class VoiceFileFinderGUI(ctk.CTk):
    def __init__(self):
        super().__init__()

        self.COLOR_BG = "#2b2024"         
        self.COLOR_PANEL = "#3d2e33"      
        self.COLOR_ACCENT = "#fdd9dd"     
        self.COLOR_ACCENT_HOVER = "#eec8cc"
        self.COLOR_BTN_TEXT = "#2b2024"
        self.COLOR_SUCCESS = "#77dd77"
        self.COLOR_ERROR = "#ff6961"

        self.title("tessell")
        self.geometry("900x650")
        self.configure(fg_color=self.COLOR_BG)
        
        self.search_path = None
        self.ser = None
        self.running = False
        self.popup_window = None
        self.command_window = None 
        self.components_ready = False
        self.vosk_model = None
        self.google_recognizer = None
        self.pending_matches = None
        self.pending_action = None 
        self.pending_target = None 
        self.undo_stack = [] 
        self.file_cache = {} 
        
        self.config = {
            'sensitivity': {
                'sound_threshold': 4000,
                'fuzzy_match_threshold': 0.3
            },
            'subjects': {
                '21st Century Literature': ['21st century', 'literature', 'philippines', 'world'],
                'DISS': ['diss', 'disciplines', 'ideas', 'social sciences'],
                'Reading and Writing': ['rws', 'reading', 'writing skills'],
                'Research Project': ['research project', 'research'],
                'Organization and Management': ['organization', 'management'],
                'UCSP': ['ucsp', 'culture', 'society', 'politics'],
                'CPAR': ['cpar', 'contemporary', 'philippine arts', 'regions'],
                'POLGOV': ['polgov', 'philippine politics', 'governance'],
                'Physical Education': ['pe', 'physical education'],
                'Applied Economics': ['applied economics', 'economics'],
                'Marketing': ['marketing'],
                'World Religions': ['iwr', 'world religions', 'belief systems'],
                'Empowerment Technologies': ['emptech', 'empowerment', 'technologies'],
                'Personal Development': ['perdev', 'personal development'],
                'Filipino': ['fpl', 'filipino', 'piling larang', 'akademik'],
                'Statistics': ['statistics', 'probability'],
                'Entrepreneurship': ['entrep', 'entrepreneurship'],
                'General Chemistry 1': ['genchem i', 'general chemistry 1', 'chemistry 1'],
                'General Chemistry 2': ['genchem ii', 'general chemistry 2', 'chemistry 2'],
                'General Math': ['genmath', 'general math', 'mathematics'],
                'Practical Research 1': ['pr i', 'practical research 1'],
                'Practical Research 2': ['pr ii', 'practical research 2'],
                'Physics': ['physics'],
                'Computer Programming': ['comprog', 'computer programming', 'programming'],
                'CESC': ['cesc', 'community engagement', 'solidarity', 'citizenship'],
                'Physical Science': ['physci', 'physical science'],
                'Front Office': ['fo', 'front office services'],
                'Housekeeping': ['hk', 'housekeeping'],
                'Accounting': ['accounting', 'accounting 1'],
                'Pagbasa at Pagsusuri': ['pagpag', 'pagbasa', 'pagsusuri'],
                'Basic Calculus': ['basic calculus', 'calculus'],
                'CSS': ['css', 'computer systems servicing'],
                'Business Finance': ['business finance', 'finance'],
                'Ethics': ['ethics'],
                'Bread & Pastry': ['b&p', 'bread', 'pastry production'],
                'Food & Beverages': ['f&b', 'food', 'beverages services'],
                'General Physics 2': ['genphysi ii', 'general physics 2', 'physics 2'],
                '3I\'S': ['3i', 'inquiries', 'investigations', 'immersion'],
                'Creative Nonfiction': ['cwf', 'creative nonfiction', 'nonfiction']
            }
        }
        
        self.grid_columnconfigure(1, weight=1)
        self.grid_rowconfigure(0, weight=1)

        self.sidebar = ctk.CTkFrame(self, width=250, corner_radius=0, fg_color=self.COLOR_PANEL)
        self.sidebar.grid(row=0, column=0, sticky="nsew")
        self.sidebar.grid_rowconfigure(6, weight=1)

        self.logo_label = ctk.CTkLabel(self.sidebar, text="tessell", 
                                     font=ctk.CTkFont(size=28, weight="bold"), text_color=self.COLOR_ACCENT)
        self.logo_label.grid(row=0, column=0, padx=20, pady=(40, 20))

        self.status_btn = ctk.CTkButton(self.sidebar, text="INITIALIZING", fg_color="#555555",
                                      text_color="white", hover=False, height=30, corner_radius=20)
        self.status_btn.grid(row=1, column=0, padx=20, pady=10)


        self.folder_btn = ctk.CTkButton(self.sidebar, text="📁 Select Folder", command=self.select_folder,
                                      fg_color=self.COLOR_ACCENT, text_color=self.COLOR_BTN_TEXT,
                                      hover_color=self.COLOR_ACCENT_HOVER, font=("Segoe UI", 14, "bold"))
        self.folder_btn.grid(row=2, column=0, padx=20, pady=20, sticky="ew")

        self.start_btn = ctk.CTkButton(self.sidebar, text="▶ START", command=self.start_system, state="disabled",
                                     fg_color=self.COLOR_SUCCESS, text_color="#1a1a1a", font=("Segoe UI", 14, "bold"))
        self.start_btn.grid(row=3, column=0, padx=20, pady=10, sticky="ew")

        self.stop_btn = ctk.CTkButton(self.sidebar, text="⏹ STOP", command=self.stop_system, state="disabled",
                                    fg_color=self.COLOR_ERROR, text_color="white", font=("Segoe UI", 14, "bold"))
        self.stop_btn.grid(row=4, column=0, padx=20, pady=10, sticky="ew")

        self.help_btn = ctk.CTkButton(self.sidebar, text="❓ Commands", command=self.show_commands,
                                    fg_color="#555555", hover_color="#666666", font=("Segoe UI", 12))
        self.help_btn.grid(row=5, column=0, padx=20, pady=20, sticky="ew")

        self.undo_btn = ctk.CTkButton(self.sidebar, text="↶ UNDO", command=self.undo_last_operation, state="disabled",
                                    fg_color="#ff9500", text_color="white", font=("Segoe UI", 12, "bold"))
        self.undo_btn.grid(row=6, column=0, padx=20, pady=10, sticky="ew")

        ctk.CTkLabel(self.sidebar, text="File Type Filter", text_color=self.COLOR_ACCENT, font=("Segoe UI", 12, "bold")).grid(row=7, column=0, padx=20, pady=(10,5))
        self.file_type_var = ctk.StringVar(value="All Files")
        self.file_type_dropdown = ctk.CTkOptionMenu(
            self.sidebar,
            variable=self.file_type_var,
            values=["All Files", "Documents (.docx)", "PDFs (.pdf)", "PowerPoint (.pptx)", "Text Files (.txt)"],
            command=self.on_filter_change,
            font=("Segoe UI", 11)
        )
        self.file_type_dropdown.grid(row=8, column=0, padx=20, pady=5, sticky="ew")

        ctk.CTkLabel(self.sidebar, text="Sensitivity", text_color=self.COLOR_ACCENT, font=("Segoe UI", 12, "bold")).grid(row=9, column=0, padx=20, pady=(10,5))
        self.sensitivity_slider = ctk.CTkSlider(self.sidebar, from_=1000, to=8000, 
                                               command=self.update_sensitivity, number_of_steps=70)  
        self.sensitivity_slider.set(self.config['sensitivity']['sound_threshold'])
        self.sensitivity_slider.grid(row=9, column=0, padx=20, pady=5, sticky="ew")
        
        self.sensitivity_label = ctk.CTkLabel(self.sidebar, text=f"{int(self.sensitivity_slider.get())}", text_color="gray70")
        self.sensitivity_label.grid(row=10, column=0, padx=20, pady=0)

        self.pc_mic_var = ctk.BooleanVar(value=True)

        self.noise_reduction_var = ctk.BooleanVar(value=False)
        self.noise_toggle = ctk.CTkCheckBox(self.sidebar, text="Noise Reduction", variable=self.noise_reduction_var, 
                                          text_color=self.COLOR_ACCENT, hover_color=self.COLOR_ACCENT_HOVER, fg_color=self.COLOR_ACCENT,
                                          font=("Segoe UI", 12))
        self.noise_toggle.grid(row=11, column=0, padx=20, pady=10, sticky="ew")

        self.info_label = ctk.CTkLabel(self.sidebar, text="Capstone Group 3 Malachi", text_color="gray60")
        self.info_label.grid(row=13, column=0, padx=20, pady=20)

        self.main_area = ctk.CTkFrame(self, fg_color="transparent")
        self.main_area.grid(row=0, column=1, sticky="nsew", padx=20, pady=20)
        self.main_area.grid_rowconfigure(1, weight=1)
        self.main_area.grid_columnconfigure(0, weight=1)

        self.console_label = ctk.CTkLabel(self.main_area, text="Activity Log", 
                                        font=("Segoe UI", 18, "bold"), text_color=self.COLOR_ACCENT)
        self.console_label.grid(row=0, column=0, sticky="w", pady=(0, 10))

        self.log_text = ctk.CTkTextbox(self.main_area, font=("Consolas", 12), text_color="#e0e0e0",
                                     fg_color="#1a1215", border_color=self.COLOR_ACCENT, border_width=1)
        self.log_text.grid(row=1, column=0, sticky="nsew")
        
        threading.Thread(target=self.initialize_components, daemon=True).start()

    def send_serial_command(self, command_bytes):
        if self.ser and self.ser.is_open:
            try:
                self.ser.write(command_bytes)
            except (OSError, serial.SerialException, Exception) as e:
                self.log(f"serial error: {e}", "WARN")
    
    def load_config(self):
        config_path = resource_path("subjects_config.json")
        default_config = {
            "subjects": {
                '21st Century Literature': ['21st century', 'literature', 'philippines', 'world'],
                'DISS': ['diss', 'disciplines', 'ideas', 'social sciences'],
                'Reading and Writing': ['rws', 'reading', 'writing skills'],
                'Research Project': ['research project', 'research'],
                'Organization and Management': ['organization', 'management'],
                'UCSP': ['ucsp', 'culture', 'society', 'politics'],
                'CPAR': ['cpar', 'contemporary', 'philippine arts', 'regions'],
                'POLGOV': ['polgov', 'philippine politics', 'governance'],
                'Physical Education': ['pe', 'physical education'],
                'Applied Economics': ['applied economics', 'economics'],
                'Marketing': ['marketing'],
                'World Religions': ['iwr', 'world religions', 'belief systems'],
                'Empowerment Technologies': ['emptech', 'empowerment', 'technologies'],
                'Personal Development': ['perdev', 'personal development'],
                'Filipino': ['fpl', 'filipino', 'piling larang', 'akademik'],
                'Statistics': ['statistics', 'probability'],
                'Entrepreneurship': ['entrep', 'entrepreneurship'],
                'General Chemistry 1': ['genchem i', 'general chemistry 1', 'chemistry 1'],
                'General Chemistry 2': ['genchem ii', 'general chemistry 2', 'chemistry 2'],
                'General Math': ['genmath', 'general math', 'mathematics'],
                'Practical Research 1': ['pr i', 'practical research 1'],
                'Practical Research 2': ['pr ii', 'practical research 2'],
                'Physics': ['physics'],
                'Computer Programming': ['comprog', 'computer programming', 'programming'],
                'CESC': ['cesc', 'community engagement', 'solidarity', 'citizenship'],
                'Physical Science': ['physci', 'physical science'],
                'Front Office': ['fo', 'front office services'],
                'Housekeeping': ['hk', 'housekeeping'],
                'Accounting': ['accounting', 'accounting 1'],
                'Pagbasa at Pagsusuri': ['pagpag', 'pagbasa', 'pagsusuri'],
                'Basic Calculus': ['basic calculus', 'calculus'],
                'CSS': ['css', 'computer systems servicing'],
                'Business Finance': ['business finance', 'finance'],
                'Ethics': ['ethics'],
                'Bread & Pastry': ['b&p', 'bread', 'pastry production'],
                'Food & Beverages': ['f&b', 'food', 'beverages services'],
                'General Physics 2': ['genphysi ii', 'general physics 2', 'physics 2'],
                '3I\'S': ['3i', 'inquiries', 'investigations', 'immersion'],
                'Creative Nonfiction': ['cwf', 'creative nonfiction', 'nonfiction']
            },
            "sensitivity": {
                "sound_threshold": 4000,
                "silence_duration": 1200,
                "fuzzy_match_threshold": 0.4
            }
        }
        try:
            if os.path.exists(config_path):
                with open(config_path, 'r') as f:
                    return json.load(f)
        except:
            pass
        return default_config

    def save_config(self):
        try:
            config_path = resource_path("subjects_config.json")
            with open(config_path, 'w') as f:
                json.dump(self.config, f, indent=2)
        except Exception as e:
            self.log(f"Config save failed: {e}", "ERROR")

    def update_sensitivity(self, value):
        threshold = int(value)
        self.config['sensitivity']['sound_threshold'] = threshold
        self.sensitivity_label.configure(text=f"{threshold}")
        self.save_config()
        if self.ser and self.ser.is_open:
            try:
                self.send_serial_command(f"THRESHOLD:{threshold}\n".encode())
            except:
                pass
    
    def on_filter_change(self, choice):
        self.log(f"Filter: {choice}", "INFO")
    
    def get_file_extension_filter(self):
        filter_map = {
            "All Files": None,
            "Documents (.docx)": [".docx", ".doc"],
            "PDFs (.pdf)": [".pdf"],
            "PowerPoint (.pptx)": [".pptx", ".ppt"],
            "Text Files (.txt)": [".txt"]
        }
        return filter_map.get(self.file_type_var.get(), None)
    
    def translate_tagalog(self, text):
        tagalog_map = {
            "hanapin": "find",
            "buksan": "open",
            "tanggalin": "delete",
            "kopyahin": "copy",
            "ilipat": "move",
            "oo": "yes",
            "hindi": "no",
            "sara": "close",
            "ayusin": "organize",
            "alphabeto": "alphabet"
        }
        
        words = text.lower().split()
        translated = []
        for word in words:
            translated.append(tagalog_map.get(word, word))
        return " ".join(translated)
    
    def show_commands(self):
        cmd_window = ctk.CTkToplevel(self)
        cmd_window.title("Voice Commands")
        cmd_window.geometry("600x500")
        
        scroll_frame = ctk.CTkScrollableFrame(cmd_window, width=560, height=460)
        scroll_frame.pack(padx=20, pady=20, fill="both", expand=True)
        
        commands_text = """
VOICE COMMANDS

📁 FILE OPERATIONS
• "Find [filename]" - Search for file
• "Delete [filename]" - Move to Deleted folder
• "Copy [filename]" - Copy path to clipboard
• "Move [filename] to [folder]" - Move file

📂 ORGANIZATION
• "Organize alphabetically" - Sort files A-Z
• "Organize by subject" - Sort by content

🇵🇭 FILIPINO/TAGALOG COMMANDS
• "Hanapin [filename]" - Find/Search
• "Buksan [filename]" - Open file
• "Tanggalin [filename]" - Delete file
• "Kopyahin [filename]" - Copy file
• "Ilipat [filename]" - Move file
• "Ayusin alphabeto" - Organize alphabetically

📋 EXAMPLES
English:
  "Find my report"
  "Delete old document"
  "Copy test file"

Filipino:
  "Hanapin test"
  "Buksan document"
  "Tanggalin old"

💡 TIPS
• Speak clearly and at normal pace
• Use exact or partial filename
• Adjust sensitivity slider if needed
• Use file type filter for faster searches
"""
        
        label = ctk.CTkLabel(scroll_frame, text=commands_text, 
                            justify="left", font=("Consolas", 12))
        label.pack(padx=10, pady=10)

    def show_calibration(self):
        cal_window = ctk.CTkToplevel(self)
        cal_window.title("Sensitivity Calibration")
        cal_window.geometry("400x300")
        cal_window.attributes("-topmost", True)
        cal_window.configure(fg_color=self.COLOR_BG)
        
        ctk.CTkLabel(cal_window, text="🎤 Calibration Mode", font=("Segoe UI", 20, "bold"), text_color=self.COLOR_ACCENT).pack(pady=15)
        ctk.CTkLabel(cal_window, text="Speak normally and watch the level.\nAdjust slider so normal speech triggers.\nBackground noise should NOT trigger.", 
                    font=("Segoe UI", 12), text_color="gray").pack(pady=10)
        
        level_label = ctk.CTkLabel(cal_window, text="Level: 0", font=("Consolas", 16, "bold"), text_color=self.COLOR_SUCCESS)
        level_label.pack(pady=20)
        
        threshold_label = ctk.CTkLabel(cal_window, text=f"Threshold: {self.config['sensitivity']['sound_threshold']}", 
                                      font=("Consolas", 14), text_color=self.COLOR_ACCENT)
        threshold_label.pack(pady=10)
        
        def update_level():
            if cal_window.winfo_exists():
                level_label.configure(text=f"Level: {getattr(self, 'last_audio_level', 0)}")
                threshold_label.configure(text=f"Threshold: {self.config['sensitivity']['sound_threshold']}")
                cal_window.after(100, update_level)
        
        update_level()

    def undo_last_operation(self):
        if not self.undo_stack:
            self.log("Nothing to undo", "WARN")
            return
        
        operation = self.undo_stack.pop()
        self.log(f"Undoing: {operation['type']}", "INFO")
        
        try:
            if operation['type'] == 'organize':
                for move in reversed(operation['moves']):
                    src = move['to']
                    dst = move['from']
                    if os.path.exists(src):
                        os.makedirs(os.path.dirname(dst), exist_ok=True)
                        shutil.move(src, dst)
                self.log(f"Restored {len(operation['moves'])} files", "SUCCESS")
            
            if not self.undo_stack:
                self.undo_btn.configure(state="disabled")
        except Exception as e:
            self.log(f"Undo failed: {e}", "ERROR")

    def log(self, message, level="INFO"):
        try:
            line_count = int(self.log_text.index('end-1c').split('.')[0])
            if line_count > 100: 
                self.log_text.delete("1.0", "21.0")
        except: pass

        timestamp = time.strftime("%H:%M:%S")
        prefix = "💬"
        if level == "ERROR": prefix = "❌"
        if level == "SUCCESS": prefix = "✅"
        if level == "WARN": prefix = "⚠️"
        
        full_msg = f"[{timestamp}] {prefix} {message}\n"
        self.log_text.insert("end", full_msg)
        self.log_text.see("end")
        
    def update_status(self, text):
        color_map = {"INITIALIZING":"#808080", "READY":self.COLOR_SUCCESS, "LISTENING":"#fdfd96",
                     "MONITORING":"#00d4ff", "CONFIRM":self.COLOR_ACCENT, "STOPPED":self.COLOR_ERROR, "ERROR":self.COLOR_ERROR}
        use_color = "#808080"
        for k,v in color_map.items(): 
            if k in text: use_color = v
        self.status_btn.configure(text=text, fg_color=use_color, text_color="#1a1a1a" if "LISTENING" in text or "READY" in text else "white")

    def show_commands(self):
        if self.command_window is None or not self.command_window.winfo_exists():
            self.command_window = ctk.CTkToplevel(self)
            self.command_window.title("Voice Commands")
            self.command_window.geometry("500x600")
            self.command_window.attributes("-topmost", True)
            self.command_window.configure(fg_color=self.COLOR_BG)
            
            ctk.CTkLabel(self.command_window, text="🎙 Command List", font=("Segoe UI", 20, "bold"), text_color=self.COLOR_ACCENT).pack(pady=15)
            
            cmds = [
                ("SEARCHING", ""),
                ("Find file name [Apple]", "Finds filename"),
                ("Find file with [Math]", "Reads inside files for text"),
                ("", ""),
                ("FILE ACTIONS", ""),
                ("Delete [filename]", "Moves file to 'Deleted' folder"),
                ("Copy [filename]", "Copies file path to clipboard"),
                ("Move [filename] to [folder]", "Moves file to specific folder"),
                ("", ""),
                ("ORGANIZING", ""),
                ("Organize Alphabetically", "Sorts files A-Z into folders"),
                ("Organize by Subject", "Sorts by content (Math, Sci...)"),
                ("", ""),
                ("SELECTION", ""),
                ("Say 'One', 'Two'...", "Selects a file from list"),
                ("Say 'Close'", "Cancels selection")
            ]
            
            scroll = ctk.CTkScrollableFrame(self.command_window, fg_color="transparent")
            scroll.pack(fill="both", expand=True, padx=10, pady=10)
            
            for cmd, desc in cmds:
                if desc == "":
                    ctk.CTkLabel(scroll, text=cmd, font=("Segoe UI", 14, "bold"), text_color="#77dd77", anchor="w").pack(fill="x", pady=(10,5))
                else:
                    f = ctk.CTkFrame(scroll, fg_color=self.COLOR_PANEL)
                    f.pack(fill="x", pady=2)
                    ctk.CTkLabel(f, text=f" 🗣 \"{cmd}\"", font=("Segoe UI", 12, "bold"), text_color="white", anchor="w").pack(fill="x", padx=10, pady=(5,0))
                    ctk.CTkLabel(f, text=desc, font=("Segoe UI", 11), text_color="gray", anchor="w").pack(fill="x", padx=10, pady=(0,5))
        else:
            self.command_window.lift()
            self.command_window.focus()

    def initialize_components(self):
        MODEL_PATH = resource_path("model")
        if not os.path.exists(MODEL_PATH):
            self.log(f"ERROR: Model missing", "ERROR")
            return
        
        try:
            self.vosk_model = Model(MODEL_PATH)
            self.google_recognizer = sr.Recognizer()
            self.log("Brains Loaded Successfully", "SUCCESS")
        except Exception as e:
            self.log(f"Error loading brains: {e}", "ERROR")
            return

        self.connect_to_esp32()

        self.update_status("✅ READY")
        self.folder_btn.configure(state="normal")
        self.components_ready = True
        self.after(500, self.select_folder)

    def find_esp32_port(self):
        ports = list(serial.tools.list_ports.comports())
        for p in ports:
            if any(x in p.description for x in ["CP210", "CH340", "USB"]): return p.device
        if ports: return ports[-1].device
        return None

    def connect_to_esp32(self):
        if self.ser and self.ser.is_open:
            try: self.ser.close()
            except: pass
            
        port = self.find_esp32_port()
        if not port:
            self.log("ESP32 Not Found. Plug it in.", "WARN")
            self.update_status("❌ NO DEVICE")
            return False
            
        try:
            self.ser = serial.Serial(port, BAUD, timeout=0.1)
            time.sleep(2) 
            self.ser.reset_input_buffer()
            self.log(f"Connected to {port}", "SUCCESS")
            return True
        except:
            self.log("Connection Failed", "ERROR")
            self.update_status("❌ ERROR")
            return False

    def select_folder(self):
        folder = filedialog.askdirectory()
        if folder:
            self.search_path = folder
            self.log(f"Target: {folder}")
            if self.components_ready: self.auto_start()

    def auto_start(self): self.after(500, self.start_system)

    def start_system(self):
        if not self.search_path: return
        
        if not self.ser or not self.ser.is_open:
            self.log("Not connected. Scanning...", "WARN")
            if not self.connect_to_esp32(): return

        try:
            self.ser.reset_input_buffer()
            self.ser.reset_output_buffer()
        except:
            self.log("Reconnecting...", "WARN")
            if not self.connect_to_esp32(): return

        self.running = True
        self.start_btn.configure(state="disabled")
        self.stop_btn.configure(state="normal")
        self.folder_btn.configure(state="disabled")
        self.update_status("🎤 MONITORING")
        threading.Thread(target=self.voice_loop, daemon=True).start()

    def stop_system(self):
        self.running = False
        self.start_btn.configure(state="normal")
        self.stop_btn.configure(state="disabled")
        self.folder_btn.configure(state="normal")
        self.update_status("⏹ STOPPED")

    def check_connection(self):
        try:
            socket.create_connection(("8.8.8.8", 53), timeout=1)
            return True
        except: return False

    def get_text_from_audio(self, wav_filename):
        if self.check_connection():
            try:
                with sr.AudioFile(wav_filename) as source:
                    audio = self.google_recognizer.record(source)
                self.google_recognizer.energy_threshold = 300
                text = self.google_recognizer.recognize_google(audio)
                self.log(f"Online: '{text}'", "SUCCESS")
                return text
            except:
                return self.get_text_vosk(wav_filename)
        else:
            return self.get_text_vosk(wav_filename)

    def get_text_vosk(self, wav_filename):
        try:
            wf = wave.open(wav_filename, "rb")
            rec = KaldiRecognizer(self.vosk_model, wf.getframerate())
            while True:
                data = wf.readframes(4000)
                if len(data) == 0: break
                rec.AcceptWaveform(data)
            res = json.loads(rec.FinalResult())
            text = res.get('text', '')
            if text: self.log(f"Offline: '{text}'", "SUCCESS")
            return text
        except: return ""

    def voice_loop(self):
        if self.pc_mic_var.get():
            if self.ser and self.ser.is_open:
                serial_thread = threading.Thread(target=self.serial_monitor, daemon=True)
                serial_thread.start()
            self.pc_mic_loop()
        else:
            self.esp32_loop()
    
    def serial_monitor(self):
        while self.running and self.pc_mic_var.get():
            try:
                if self.ser and self.ser.in_waiting > 0:
                    raw = self.ser.read(self.ser.in_waiting)
                    decoded = raw.decode('utf-8', 'ignore')
                    
                    if "THRESHOLD:" in decoded:
                        try:
                            threshold_str = decoded.split("THRESHOLD:")[1].split()[0]
                            threshold = int(threshold_str)
                            
                            def update_ui(t):
                                try:
                                    self.sensitivity_slider.set(t)
                                    self.sensitivity_label.configure(text=f"{t}")
                                except:
                                    pass
                            
                            self.after(0, lambda: update_ui(threshold))
                            
                        except:
                            pass
                                
                time.sleep(0.1)
                
            except:
                break
    
    def pc_mic_loop(self):
        waiting_for_confirmation = False
        self.log("PC Microphone Mode Active", "SUCCESS")
        
        while self.running:
            try:
                self.update_status("🎤 PC MIC - Ready")
                
                with sr.Microphone() as source:
                    if self.noise_reduction_var.get():
                        self.log("Adjusting for noise...", "INFO")
                        self.google_recognizer.adjust_for_ambient_noise(source, duration=1.0)
                        self.google_recognizer.dynamic_energy_threshold = True
                        self.google_recognizer.energy_threshold = 4000 
                    else:
                        self.google_recognizer.adjust_for_ambient_noise(source, duration=0.5)
                        self.google_recognizer.dynamic_energy_threshold = True

                    self.update_status("🎙 LISTENING...")
                    audio = self.google_recognizer.listen(source, timeout=10, phrase_time_limit=5)
                
                self.update_status("⏳ Processing...")
                
                try:
                    if self.check_connection():
                        text = self.google_recognizer.recognize_google(audio)
                        self.log(f"Heard: '{text}'", "SUCCESS")
                    else:
                        with open(TEMP_WAV, "wb") as f:
                            f.write(audio.get_wav_data())
                        text = self.get_text_vosk(TEMP_WAV)
                    
                    if text:
                        text = text.lower()
                        
                        if waiting_for_confirmation and self.pending_matches:
                            if any(w in text for w in ["close", "cancel", "no", "stop"]):
                                self.log("Cancelled.")
                                self.close_popup()
                                self.pending_matches = None
                                self.pending_action = None
                                self.pending_target = None
                                waiting_for_confirmation = False
                            else:
                                choice = text_to_int(text)
                                if len(self.pending_matches) == 1 and any(w in text for w in ["yes", "open", "sure", "ok"]):
                                    choice = 1
                                if choice and 1 <= choice <= len(self.pending_matches):
                                    self.execute_pending_action(choice-1)
                                    waiting_for_confirmation = False
                                else:
                                    self.log(f"Say a number 1-{len(self.pending_matches)}", "WARN")
                        else:
                            self.process_command(text)
                            if self.pending_matches:
                                waiting_for_confirmation = True
                    else:
                        self.log("Didn't catch that", "WARN")
                        
                except sr.UnknownValueError:
                    self.log("Didn't understand", "WARN")
                except sr.RequestError:
                    self.log("Connection error", "ERROR")
                except Exception as e:
                    self.log(f"Error: {e}", "ERROR")
                    
            except sr.WaitTimeoutError:
                continue
            except Exception as e:
                self.log(f"Mic error: {e}", "ERROR")
                time.sleep(1)
    
    def esp32_loop(self):
        buffer = bytearray()
        recording = False
        waiting_for_confirmation = False
        
        while self.running:
            try:
                try:
                    if self.ser.in_waiting > 0:
                        raw = self.ser.read(self.ser.in_waiting)
                    else:
                        time.sleep(0.01)
                        continue
                except (OSError, serial.SerialException):
                    self.log("⚠️ Device Disconnected!", "ERROR")
                    self.update_status("❌ DISCONNECTED")
                    self.stop_system()
                    self.ser = None
                    break 

                decoded = raw.decode('utf-8', 'ignore')
                if "THRESHOLD:" in decoded:
                    try:
                        threshold_str = decoded.split("THRESHOLD:")[1].split()[0]
                        threshold = int(threshold_str)
                        self.sensitivity_slider.set(threshold) 
                        self.sensitivity_label.configure(text=f"{threshold}") 
                        self.log(f"🎛️ Pot: {threshold}", "INFO")
                    except:
                        pass

                if not recording and "START_REC" in decoded:
                    recording = True
                    buffer = bytearray()
                    self.update_status("🎙 LISTENING")
                
                if recording:
                    buffer.extend(raw)
                    if b"STOP_REC" in buffer:
                        recording = False
                        clean = buffer.split(b"STOP_REC")[0]
                        with wave.open(TEMP_WAV, 'wb') as wf:
                            wf.setnchannels(1); wf.setsampwidth(2); wf.setframerate(16000)
                            wf.writeframes(clean)
                        
                        try:
                            text = self.get_text_from_audio(TEMP_WAV)
                            if text: 
                                text = text.lower()
                                
                                if waiting_for_confirmation and self.pending_matches:
                                    if any(w in text for w in ["close", "cancel", "no", "stop"]):
                                        self.log("Cancelled.")
                                        self.close_popup()
                                        self.pending_matches = None
                                        self.pending_action = None
                                        self.pending_target = None
                                        waiting_for_confirmation = False
                                        self.update_status("🎤 MONITORING")
                                    else:
                                        choice = text_to_int(text)
                                        if len(self.pending_matches) == 1 and any(w in text for w in ["yes", "open", "sure", "ok"]): choice = 1
                                        if choice and 1 <= choice <= len(self.pending_matches):
                                            self.execute_pending_action(choice-1)
                                            self.send_serial_command(b"FOUND\n") 
                                            waiting_for_confirmation = False
                                            self.update_status("🎤 MONITORING")
                                        else:
                                            self.log(f"Say a number 1-{len(self.pending_matches)}", "WARN")

                                else:
                                    self.process_command(text)
                                    if self.pending_matches: 
                                        waiting_for_confirmation = True

                            else: self.update_status("🎤 MONITORING")
                        except Exception as e:
                            self.log(f"Logic Error: {e}", "ERROR")
                        buffer = bytearray()
            except: break

    def process_command(self, text):
        text = self.translate_tagalog(text)
        
        if "arrange" in text or "organize" in text:
            if "alphabet" in text: self.organize_alphabetically()
            elif "subject" in text or "content" in text: self.organize_by_subject()
            else: self.log("Say 'Alphabetically' or 'By Subject'", "WARN")
        
        elif "delete" in text or "remove" in text:
            clean = text.replace("delete", "").replace("remove", "").replace("file", "").strip()
            if clean:
                self.delete_file(clean)
            else:
                self.log("Say filename to delete", "WARN")
        
        elif "copy" in text:
            clean = text.replace("copy", "").replace("file", "").strip()
            if clean:
                self.copy_file(clean)
            else:
                self.log("Say filename to copy", "WARN")
        
        elif "move" in text:
            if "to" in text:
                parts = text.split("to")
                filename = parts[0].replace("move", "").replace("file", "").strip()
                folder = parts[1].strip() if len(parts) > 1 else ""
                if filename and folder:
                    self.move_file(filename, folder)
                else:
                    self.log("Say 'move [file] to [folder]'", "WARN")
            else:
                self.log("Say 'move [file] to [folder]'", "WARN")
        
        elif "find" in text or "search" in text:
            clean = text.replace("search", "").replace("find", "").replace("for", "").strip()
            if not clean:
                self.log("Say filename (e.g. 'Search for Notes')", "WARN")
                return 
            if "with" in clean:
                parts = clean.split("with")
                if len(parts) > 1 and parts[1].strip(): self.find_by_content(parts[1].strip())
                else: self.log("Say a word after 'with'", "WARN")
            else:
                keyword = clean.replace("file name", "").strip()
                if keyword: self.find_by_name(keyword)
                else: self.log("Say the filename", "WARN")
        else: self.update_status("🎤 MONITORING")

    def organize_alphabetically(self):
        self.log("Sorting files A-Z...", "INFO")
        self.send_serial_command(b"SEARCHING\n")
        count = 0
        moves = []
        try:
            for file in os.listdir(self.search_path):
                if os.path.isfile(os.path.join(self.search_path, file)):
                    first = file[0].upper()
                    if not first.isalpha(): first = "#"
                    target = os.path.join(self.search_path, first)
                    if not os.path.exists(target): os.makedirs(target)
                    
                    src = os.path.join(self.search_path, file)
                    dst = os.path.join(target, file)
                    shutil.move(src, dst)
                    moves.append({'from': src, 'to': dst})
                    count += 1
            
            self.undo_stack.append({'type': 'organize', 'moves': moves})
            self.undo_btn.configure(state="normal")
            
            self.log(f"Moved {count} files.", "SUCCESS")
            self.send_serial_command(b"FOUND\n")
        except Exception as e:
            self.log(f"Error organizing: {e}", "ERROR")
            self.send_serial_command(b"NOTFOUND\n")

    def organize_by_subject(self):
        self.log("Organizing by Subject...", "INFO")
        self.send_serial_command(b"SEARCHING\n")
        subjects = self.config['subjects']
        count = 0
        moves = []
        try:
            for file in os.listdir(self.search_path):
                file_path = os.path.join(self.search_path, file)
                if os.path.isfile(file_path):
                    content = get_file_content(file_path)
                    
                    for subject, keywords in subjects.items():
                        if any(k.lower() in content for k in keywords) or any(k.lower() in file.lower() for k in keywords):
                            target = os.path.join(self.search_path, subject)
                            if not os.path.exists(target): os.makedirs(target)
                            
                            dst = os.path.join(target, file)
                            shutil.move(file_path, dst)
                            moves.append({'from': file_path, 'to': dst})
                            count += 1
                            break
            
            self.undo_stack.append({'type': 'organize', 'moves': moves})
            self.undo_btn.configure(state="normal")
            
            self.log(f"Organized {count} files.", "SUCCESS")
            self.send_serial_command(b"FOUND\n")
        except Exception as e:
            self.log(f"Error: {e}", "ERROR")
            self.send_serial_command(b"NOTFOUND\n")
    
    def delete_file(self, keyword):
        self.log(f"Deleting '{keyword}'...", "INFO")
        self.find_by_name(keyword)
        if self.pending_matches:
            self.log("Delete this file?", "CONFIRM")
            self.pending_action = "delete"
            self.send_serial_command(b"FOUND\n")
    
    def copy_file(self, keyword):
        self.log(f"Copying '{keyword}'...", "INFO")
        self.find_by_name(keyword)
        if self.pending_matches:
            self.log("Copy this file?", "CONFIRM")
            self.pending_action = "copy"
            self.send_serial_command(b"FOUND\n")

    def move_file(self, keyword, target_folder_name):
        self.log(f"Moving '{keyword}'...", "INFO")
        self.find_by_name(keyword)
        if self.pending_matches:
            target_path = None
            for root, dirs, _ in os.walk(self.search_path):
                for d in dirs:
                    if target_folder_name.lower() in d.lower():
                        target_path = os.path.join(root, d)
                        break
            
            if target_path:
                self.pending_target = target_path
                self.log(f"Move to {os.path.basename(target_path)}?", "CONFIRM")
                self.pending_action = "move"
                self.send_serial_command(b"FOUND\n")
            else:
                self.log(f"Folder '{target_folder_name}' not found", "WARN")
                self.pending_matches = None
                self.send_serial_command(b"NOTFOUND\n")

    def find_by_content(self, keyword):
        self.log(f"Reading content for '{keyword}'...", "INFO")
        self.send_serial_command(b"SEARCHING\n")
        matches = []
        for root, dirs, files in os.walk(self.search_path):
            for file in files:
                ext_filter = self.get_file_extension_filter()
                if ext_filter and not any(file.lower().endswith(e) for e in ext_filter):
                    continue
                    
                path = os.path.join(root, file)
                content = get_file_content(path)
                if keyword.lower() in content:
                    matches.append({'name': file, 'path': path, 'score': 100})
        
        self.handle_results(matches)

    def find_by_name(self, filename):
        self.log(f"Searching: '{filename}'", "INFO")
        self.send_serial_command(b"SEARCHING\n")
        matches = []
        threshold = self.config['sensitivity']['fuzzy_match_threshold']
        for root, dirs, files in os.walk(self.search_path):
            for file in files:
                ext_filter = self.get_file_extension_filter()
                if ext_filter and not any(file.lower().endswith(e) for e in ext_filter):
                    continue
                
                ratio = difflib.SequenceMatcher(None, filename.lower(), file.lower()).ratio()
                if ratio > threshold:
                    matches.append({'name': file, 'path': os.path.join(root, file), 'score': ratio})
        
        self.handle_results(matches)

    def handle_results(self, matches):
        matches.sort(key=lambda x: x['score'], reverse=True)
        if matches:
            self.pending_matches = matches
            self.pending_action = "open" 
            self.send_serial_command(b"FOUND\n")
            self.after(0, lambda: self.show_modern_popup(matches))
        else:
            self.send_serial_command(b"NOTFOUND\n")
            self.log("No matches found.", "WARN")
            self.pending_matches = None
            self.pending_action = None

    def show_matches(self, matches, action="open"):
        self.pending_action = action
        self.pending_matches = matches
        self.send_serial_command(b"FOUND\n")
        self.after(0, lambda: self.show_modern_popup(matches))
        
    def execute_pending_action(self, index):
        if not self.pending_matches or index >= len(self.pending_matches): return
        
        match = self.pending_matches[index]
        filepath = match['path']
        filename = match['name']
        
        try:
            if self.pending_action == 'open':
                os.startfile(filepath)
                self.log(f"Opened: {filename}", "SUCCESS")
                
            elif self.pending_action == 'delete':
                try:
                    from send2trash import send2trash
                    send2trash(filepath)
                    self.log(f"Recycled: {filename}", "SUCCESS")
                except ImportError:
                    trash_path = os.path.join(self.search_path, "Deleted")
                    if not os.path.exists(trash_path): os.makedirs(trash_path)
                    shutil.move(filepath, os.path.join(trash_path, filename))
                    self.log(f"Moved to Deleted folder: {filename}", "SUCCESS")

            elif self.pending_action == 'copy':
                try:
                    copy_files_to_clipboard([filepath])
                    self.log(f"File copied: {filename}", "SUCCESS")
                except Exception as e:
                    self.log(f"Copy failed: {e}", "ERROR")

            elif self.pending_action == 'move':
                if self.pending_target:
                    shutil.move(filepath, os.path.join(self.pending_target, filename))
                    self.log(f"Moved to {os.path.basename(self.pending_target)}", "SUCCESS")
            
            self.pending_matches = None
            self.pending_action = None
            self.pending_target = None
            self.close_popup()
            
        except Exception as e:
            self.log(f"Action failed: {e}", "ERROR")

    def show_modern_popup(self, items):
        if self.popup_window: 
            try: self.popup_window.destroy()
            except: pass
            
        self.popup_window = ctk.CTkToplevel(self)
        self.popup_window.title("Results")
        self.popup_window.geometry("800x500") 
        self.popup_window.attributes("-topmost", True)
        self.popup_window.configure(fg_color=self.COLOR_BG)
        
        container = ctk.CTkFrame(self.popup_window, fg_color="transparent")
        container.pack(fill="both", expand=True, padx=10, pady=10)
        
        left_panel = ctk.CTkFrame(container, fg_color="transparent")
        left_panel.pack(side="left", fill="both", expand=True, padx=(0, 10))
        
        ctk.CTkLabel(left_panel, text=f"📂 Found {len(items)} Files", font=("Segoe UI", 20, "bold"), text_color=self.COLOR_SUCCESS).pack(pady=10)
        
        scroll_frame = ctk.CTkScrollableFrame(left_panel, fg_color="transparent")
        scroll_frame.pack(fill="both", expand=True)
        
        right_panel = ctk.CTkFrame(container, width=300, fg_color=self.COLOR_PANEL)
        right_panel.pack(side="right", fill="both", padx=(10, 0))
        
        ctk.CTkLabel(right_panel, text="👁 Preview", font=("Segoe UI", 16, "bold"), text_color=self.COLOR_ACCENT).pack(pady=10)
        preview_text = ctk.CTkTextbox(right_panel, font=("Consolas", 11), text_color="#e0e0e0", fg_color="#1a1215")
        preview_text.pack(fill="both", expand=True, padx=10, pady=10)
        
        def show_preview(path):
            content = get_file_content(path)
            preview_text.delete("1.0", "end")
            preview_text.insert("1.0", content[:1000] if content else "no preview available")
        
        for i, match in enumerate(items):
            btn_text = f"{i+1}. {match['name']}"
            card = ctk.CTkButton(scroll_frame, text=btn_text, anchor="w", fg_color="#4a3b40", 
                               hover_color="#5c4a50", corner_radius=10, height=40, 
                               font=("Segoe UI", 12), text_color="white",
                               command=lambda p=match['path']: show_preview(p))
            card.pack(fill="x", pady=5, padx=5)

        ctk.CTkLabel(left_panel, text="🎤 Say 'ONE', 'OPEN ONE'...", font=("Segoe UI", 14), text_color=self.COLOR_ACCENT).pack(pady=15)
        self.update_status("❓ CONFIRM")
        
        if items: show_preview(items[0]['path'])
        
    def close_popup(self):
        if self.popup_window:
            self.popup_window.destroy()
            self.popup_window = None

if __name__ == "__main__":
    app = VoiceFileFinderGUI()
    app.mainloop()
